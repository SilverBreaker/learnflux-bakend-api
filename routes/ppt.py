import uuid
import os
from fastapi import APIRouter, Depends, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from auth_utils import get_current_user
from config import UPLOAD_DIR, OUTPUT_DIR
from services.pdf_extractor import extract_text_from_pdf

router = APIRouter()

def _generate_ppt(title: str, text: str, out_path: str) -> int:
    """Generate .pptx from PDF text — no AI needed."""

    # Split text into slide chunks
    words = text.split()
    chunks = []
    chunk_size = 80
    for i in range(0, min(len(words), 640), chunk_size):
        chunks.append(" ".join(words[i:i + chunk_size]))

    slide_titles = [
        "Introduction", "Key Concepts", "Main Topics",
        "Important Points", "Details", "Analysis",
        "Summary", "Conclusion"
    ]

    slides_data = []
    for i, chunk in enumerate(chunks[:8]):
        sentences = chunk.split('.')
        bullets = [s.strip() for s in sentences if len(s.strip()) > 10][:4]
        slides_data.append({
            "title": slide_titles[i] if i < len(slide_titles) else f"Slide {i+1}",
            "bullets": bullets if bullets else [chunk[:100]]
        })

    if not slides_data:
        slides_data = [{"title": title, "bullets": ["No content could be extracted"]}]

    # Build PPTX
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide_layout = prs.slide_layouts[1]

    for slide_info in slides_data:
        slide = prs.slides.add_slide(slide_layout)

        # Title
        tf = slide.shapes.title.text_frame
        tf.text = slide_info.get("title", "")
        tf.paragraphs[0].runs[0].font.size = Pt(32)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x9B, 0x6B, 0xFF)

        # Bullets
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(slide_info.get("bullets", [])):
            p = body.add_paragraph() if i > 0 else body.paragraphs[0]
            p.text = bullet
            p.level = 0
            if p.runs:
                p.runs[0].font.size = Pt(18)

    prs.save(out_path)
    return len(slides_data)


# POST /ppt/generate
@router.post("/generate")
async def generate_ppt(
    file: UploadFile = File(...),
    user_id: str = Depends(get_current_user),
):
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files allowed")

    job_id   = str(uuid.uuid4())
    file_path = os.path.join(UPLOAD_DIR, f"ppt_{job_id}.pdf")
    out_path  = os.path.join(OUTPUT_DIR, f"ppt_{job_id}.pptx")

    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)

    text        = extract_text_from_pdf(file_path)
    slide_count = _generate_ppt(file.filename.replace(".pdf", ""), text, out_path)

    # Clean up uploaded pdf
    if os.path.exists(file_path):
        os.remove(file_path)

    return {
        "ppt_url":     f"/ppt/download/{job_id}",
        "slide_count": slide_count,
        "title":       file.filename.replace(".pdf", "")
    }


# GET /ppt/download/{job_id}
@router.get("/download/{job_id}")
async def download_ppt(job_id: str):
    path = os.path.join(OUTPUT_DIR, f"ppt_{job_id}.pptx")
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="PPT not found")
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="presentation.pptx"
    )